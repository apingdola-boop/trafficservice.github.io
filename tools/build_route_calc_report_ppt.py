# -*- coding: utf-8 -*-
"""One-off: ROUTE CALC 개선 보고용 PPT 생성 (python-pptx)"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN


def _body(tf, lines):
    tf.clear()
    p = tf.paragraphs[0]
    for i, line in enumerate(lines):
        if i > 0:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(18 if i == 0 and len(lines) > 3 else 16)
        p.level = 0


def main():
    root = Path(__file__).resolve().parents[1]
    # 영문 파일명: 터미널·일부 도구에서 한글 경로 깨짐 방지
    out = root / "ROUTE_CALC_report.pptx"

    prs = Presentation()
    prs.slide_width = prs.slide_width  # default 16:9

    # 1. Title
    s0 = prs.slides.add_slide(prs.slide_layouts[0])
    s0.shapes.title.text = "ROUTE CALC 개선 보고"
    st = s0.placeholders[1].text_frame
    st.text = "실경로·유류비 계산 / 엑셀·H-챗 연동 / 검색 UI"

    # 2. 개요
    s1 = prs.slides.add_slide(prs.slide_layouts[1])
    s1.shapes.title.text = "개요"
    _body(
        s1.shapes.placeholders[1].text_frame,
        [
            "자동 검색 결과를 정산표·H-챗에 맞게 표기 통일",
            "면접원·프로젝트 검색 탭 배치 및 모바일 최적화",
            "수동 검색 계산 로직은 변경 없음",
        ],
    )

    # 3. 자동 검색 → 엑셀·H-챗
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "자동 검색 → 엑셀·H-챗"
    _body(
        s2.shapes.placeholders[1].text_frame,
        [
            "주소: 가능 시 「OO구 OO동」 축약",
            "경유: 동당 한 칸, 연속 동일 동 병합 (부수만큼 행 반복 없음)",
            "거리: 반올림 총 km를 구간 수로 나눈 정수 (합계 일치)",
            "내비 원시 구간은 내부용으로만 보관, 전송·엑셀은 합성 구간 사용",
        ],
    )

    # 4. UI 데스크톱
    s3 = prs.slides.add_slide(prs.slide_layouts[1])
    s3.shapes.title.text = "검색 UI (데스크톱)"
    _body(
        s3.shapes.placeholders[1].text_frame,
        [
            "면접원·프로젝트 탭을 왼쪽에 세로로 통합",
            "프로젝트 패널도 왼쪽 슬라이드",
            "한쪽 열면 다른 쪽 자동 닫힘 (겹침 방지)",
        ],
    )

    # 5. 모바일
    s4 = prs.slides.add_slide(prs.slide_layouts[1])
    s4.shapes.title.text = "검색 UI (모바일)"
    _body(
        s4.shapes.placeholders[1].text_frame,
        [
            "하단 중앙 가로 pill 버튼 (본문 가리지 않음)",
            "면접원·프로젝트 동일 스타일 (CSS 순서 이슈 수정)",
            "하단 안전영역·본문 하단 여백 반영",
        ],
    )

    # 6. 범위
    s5 = prs.slides.add_slide(prs.slide_layouts[1])
    s5.shapes.title.text = "적용 범위·참고"
    _body(
        s5.shapes.placeholders[1].text_frame,
        [
            "주요 변경 파일: index.html (GitHub Pages 배포)",
            "엑셀 업로드·H-챗 postMessage 요약 문구 동일 규칙",
            "보고·발표 시 스크린샷: 자동 검색 정산표 + 모바일 하단 탭",
        ],
    )

    prs.save(out)
    print("Wrote:", out)


if __name__ == "__main__":
    main()
